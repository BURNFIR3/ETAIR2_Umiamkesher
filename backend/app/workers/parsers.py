"""
File family-specific parsers.
Each parser returns:
{
    "title": str,
    "language": str,
    "family_data": dict,
    "chunks": [{"content", "chunk_type", "page_number", ...}],
    "entities": [{"type", "value", "confidence"}],
}
"""
import io
import re
from typing import Any, Dict, List, Optional


EQUIPMENT_PATTERN = re.compile(r'\b([A-Z]{1,3}-?\d{3,6}[A-Z]?)\b')
STANDARD_PATTERN = re.compile(r'\b(ISO|IEC|ASME|ASTM|ANSI|API|BS|DIN|EN)\s?[\d\-]+\b')


def _strip_null_bytes(data: Any) -> Any:
    if isinstance(data, str):
        return data.replace("\x00", "")
    elif isinstance(data, dict):
        return {k: _strip_null_bytes(v) for k, v in data.items()}
    elif isinstance(data, list):
        return [_strip_null_bytes(item) for item in data]
    return data


def parse_file(content: bytes, filename: str, mime_type: str, file_family: str) -> Dict[str, Any]:
    """Route to the correct parser based on file family and sanitize null bytes."""
    parsers = {
        "text_office": _parse_text_office,
        "table": _parse_table,
        "image": _parse_image,
        "audio": _parse_audio,
        "cad": _parse_cad,
        "operational": _parse_operational,
    }
    parser = parsers.get(file_family, _parse_unknown)
    result = parser(content, filename, mime_type)
    return _strip_null_bytes(result)


# ─── Text / Office ────────────────────────────────────────────────────────────

def _parse_text_office(content: bytes, filename: str, mime_type: str) -> Dict:
    ext = filename.lower().rsplit(".", 1)[-1]

    if ext == "pdf":
        return _parse_pdf(content, filename)
    elif ext in ("docx", "doc"):
        return _parse_docx(content, filename)
    elif ext in ("pptx", "ppt"):
        return _parse_pptx(content, filename)
    elif ext in ("txt", "md", "html", "htm"):
        return _parse_plaintext(content, filename)
    else:
        return _parse_plaintext(content, filename)


def _parse_pdf(content: bytes, filename: str) -> Dict:
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(stream=content, filetype="pdf")
        chunks = []
        all_text = []
        is_scanned = True

        # First pass: check if document has embedded text
        doc_text_total = sum(len(page.get_text("text").strip()) for page in doc)
        if doc_text_total > 100:
            is_scanned = False

        max_ocr_pages = 10  # Cap OCR pages to prevent OOM / infinite timeouts on long scanned PDFs
        ocr_count = 0

        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text").strip()

            if is_scanned and len(text) < 15 and ocr_count < max_ocr_pages:
                try:
                    pix = page.get_pixmap(matrix=fitz.Matrix(1.0, 1.0))
                    png_bytes = pix.tobytes("png")
                    pix = None  # Free pixmap buffer immediately
                    ocr_result = _ocr_image_bytes(png_bytes, "png") or ""
                    if ocr_result.strip():
                        text = ocr_result.strip()
                    ocr_count += 1
                except Exception:
                    pass

            if text:
                all_text.append(text)
                for chunk_text in _split_into_chunks(text):
                    chunks.append({
                        "content": chunk_text,
                        "chunk_type": "page",
                        "page_number": page_num + 1,
                    })

        page_cnt = len(doc)
        doc.close()
        full_text = "\n".join(all_text)
        if not full_text.strip():
            full_text = f"[Scanned/unparseable PDF document: {filename}]"
            chunks.append({"content": full_text, "chunk_type": "page", "page_number": 1})

        entities = _extract_entities(full_text)

        return {
            "title": _guess_title(filename, full_text),
            "language": "en",
            "family_data": {
                "page_count": page_cnt,
                "is_scanned": is_scanned,
                "word_count": len(full_text.split()),
            },
            "chunks": _group_into_sections(chunks, max_chars=2200),
            "entities": entities,
            "keywords": _extract_keywords(full_text, top_n=20),
        }
    except Exception as e:
        return _error_result(filename, str(e))


def _parse_docx(content: bytes, filename: str) -> Dict:
    try:
        from docx import Document

        doc = Document(io.BytesIO(content))
        chunks = []
        all_text = []
        current_heading = None

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            all_text.append(text)
            style = para.style.name if para.style else ""

            if "Heading" in style:
                current_heading = text
                chunks.append({"content": text, "chunk_type": "heading"})
            else:
                chunks.append({"content": text, "chunk_type": "paragraph"})

        # Also extract table content
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    all_text.append(row_text)
                    chunks.append({"content": row_text, "chunk_type": "table_row"})

        full_text = "\n".join(all_text)
        return {
            "title": _guess_title(filename, full_text),
            "language": "en",
            "family_data": {
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
            },
            "chunks": _group_into_sections(_merge_small_chunks(chunks), max_chars=2200),
            "entities": _extract_entities(full_text),
            "keywords": _extract_keywords(full_text, top_n=20),
        }
    except Exception as e:
        return _error_result(filename, str(e))


def _parse_pptx(content: bytes, filename: str) -> Dict:
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        chunks = []
        all_text = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                notes = slide.notes_slide.notes_text_frame.text.strip()
                slide_texts.append(f"[Notes] {notes}")

            if slide_texts:
                slide_content = "\n".join(slide_texts)
                all_text.append(slide_content)
                chunks.append({
                    "content": slide_content,
                    "chunk_type": "slide",
                    "slide_number": slide_num,
                })

        full_text = "\n".join(all_text)
        return {
            "title": _guess_title(filename, full_text),
            "language": "en",
            "family_data": {"slide_count": len(prs.slides)},
            "chunks": _group_into_sections(chunks, max_chars=2200),
            "entities": _extract_entities(full_text),
            "keywords": _extract_keywords(full_text, top_n=20),
        }
    except Exception as e:
        return _error_result(filename, str(e))


def _parse_plaintext(content: bytes, filename: str) -> Dict:
    try:
        text = content.decode("utf-8", errors="replace")
        chunks = [{"content": c, "chunk_type": "paragraph"} for c in _split_into_chunks(text)]
        return {
            "title": filename,
            "language": "en",
            "family_data": {"char_count": len(text)},
            "chunks": _group_into_sections(chunks, max_chars=2200),
            "entities": _extract_entities(text),
            "keywords": _extract_keywords(text, top_n=20),
        }
    except Exception as e:
        return _error_result(filename, str(e))


# ─── Tables / Structured ──────────────────────────────────────────────────────

def _parse_table(content: bytes, filename: str, mime_type: str) -> Dict:
    try:
        import pandas as pd

        ext = filename.lower().rsplit(".", 1)[-1]
        if ext in ("xlsx", "xls", "ods"):
            xls = pd.ExcelFile(io.BytesIO(content))
            sheet_names = xls.sheet_names
            chunks = []
            all_headers = {}

            for sheet in sheet_names:
                df = xls.parse(sheet)
                headers = list(df.columns.astype(str))
                all_headers[sheet] = headers
                # Chunk in row groups of 50
                for start in range(0, len(df), 50):
                    end = min(start + 50, len(df))
                    subset = df.iloc[start:end].to_string(index=False)
                    chunks.append({
                        "content": f"Sheet: {sheet}\nColumns: {', '.join(headers)}\n{subset}",
                        "chunk_type": "row_group",
                        "row_start": start,
                        "row_end": end,
                    })

            full_text_xls = " ".join(str(sheet_names)) + " " + " ".join(
                h for hs in all_headers.values() for h in hs
            )
            return {
                "title": filename,
                "language": "en",
                "family_data": {"sheet_names": sheet_names, "column_headers": all_headers},
                "chunks": _group_into_sections(chunks, max_chars=2200),
                "entities": [],
                "keywords": _extract_keywords(full_text_xls, top_n=20),
            }

        elif ext in ("csv", "tsv"):
            sep = "\t" if ext == "tsv" else ","
            df = pd.read_csv(io.BytesIO(content), sep=sep, on_bad_lines="skip")
            headers = list(df.columns.astype(str))
            chunks = []
            for start in range(0, len(df), 50):
                end = min(start + 50, len(df))
                subset = df.iloc[start:end].to_string(index=False)
                chunks.append({
                    "content": f"Columns: {', '.join(headers)}\n{subset}",
                    "chunk_type": "row_group",
                    "row_start": start,
                    "row_end": end,
                })
            return {
                "title": filename,
                "language": "en",
                "family_data": {"row_count": len(df), "column_headers": headers},
                "chunks": _group_into_sections(chunks, max_chars=2200),
                "entities": [],
                "keywords": _extract_keywords(" ".join(headers), top_n=20),
            }

    except Exception as e:
        return _error_result(filename, str(e))


# ─── Images / Scans ───────────────────────────────────────────────────────────

def _parse_image(content: bytes, filename: str, mime_type: str) -> Dict:
    try:
        ext = filename.lower().rsplit(".", 1)[-1]
        ocr_text = _ocr_image_bytes(content, ext) or ""
        entities = _extract_entities(ocr_text, is_image=True)
        keywords = _extract_keywords(ocr_text, top_n=30)

        # Ensure image has multiple rich entities: supplement regex entities with concepts & keywords
        seen_vals = {e["value"].lower() for e in entities}
        for kw in keywords:
            if kw.lower() not in seen_vals and len(kw) > 3:
                seen_vals.add(kw.lower())
                entities.append({"type": "concept", "value": kw, "confidence": 0.85})
                if len(entities) >= 35:
                    break

        chunks = []
        if ocr_text:
            for chunk in _split_into_chunks(ocr_text):
                chunks.append({"content": chunk, "chunk_type": "region"})
        else:
            chunks.append({"content": f"Image diagram scan: {filename}", "chunk_type": "region"})

        return {
            "title": filename,
            "language": "en",
            "family_data": {"ocr_word_count": len(ocr_text.split()), "has_ocr_text": bool(ocr_text), "entity_count": len(entities)},
            "chunks": chunks,
            "entities": entities,
            "keywords": keywords[:20],
        }
    except Exception as e:
        return _error_result(filename, str(e))


def _ocr_image_bytes(content: bytes, ext: str = "png") -> Optional[str]:
    try:
        import pytesseract
        from PIL import Image, ImageEnhance
        img = Image.open(io.BytesIO(content)).convert("RGB")
        # Upscale 2x and enhance contrast/sharpness for handwriting detection
        w, h = img.size
        img = img.resize((max(w * 2, 800), max(h * 2, 800)), Image.Resampling.LANCZOS)
        img = img.convert("L")
        img = ImageEnhance.Contrast(img).enhance(2.2)
        img = ImageEnhance.Sharpness(img).enhance(2.0)

        results = []
        for psm in [3, 6, 11]:
            try:
                txt = pytesseract.image_to_string(img, config=f"--psm {psm}", timeout=10).strip()
                if txt:
                    results.append(txt)
            except Exception:
                pass
        img.close()

        combined_lines = []
        seen_lines = set()
        for res_text in results:
            for line in res_text.splitlines():
                cl = line.strip()
                if cl and cl.lower() not in seen_lines:
                    seen_lines.add(cl.lower())
                    combined_lines.append(cl)
        return "\n".join(combined_lines) if combined_lines else None
    except Exception:
        return None


# ─── Audio ───────────────────────────────────────────────────────────────────

def _parse_audio(content: bytes, filename: str, mime_type: str) -> Dict:
    """
    Audio parsing pipeline — optimised for minimum storage, maximum query clarity.

    Flow:
    1. Write bytes to a temp file (Whisper API requires a file-like with a name)
    2. Call OpenAI Whisper with verbose_json + segment timestamps
    3. Merge short segments into ~AUDIO_CHUNK_MAX_TOKENS chunks (better LLM context)
    4. Store the compressed full transcript in family_data (metadata search works
       without needing an embedding lookup)
    5. Optionally delete raw audio from MinIO (KEEP_AUDIO_RAW=False by default)
       — the transcript is fully preserved in the DB, so nothing is lost.

    Storage comparison:
      Raw audio:            100MB+  per file (if kept)
      Transcript segments:  ~5KB    in file_chunks (text only)
      Compressed transcript:~2KB    in family_data JSONB
      Embeddings:           ~100KB  per file (merged chunks → fewer vectors)
    """
    try:
        import base64
        import zlib
        import tempfile
        import os
        import math
        from app.config import settings

        if not settings.GROQ_API_KEY:
            return _error_result(filename, "GROQ_API_KEY not set for audio transcription")

        from groq import Groq
        client = Groq(api_key=settings.GROQ_API_KEY)
        ext = filename.lower().rsplit(".", 1)[-1]

        # ── Step 1: Write to temp file ──────────────────────────────────────
        with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        # ── Step 2: Transcribe with Groq Whisper (whisper-large-v3) ─────────
        # Groq's Whisper is free, fast, and accurate.
        # verbose_json gives us segments with timestamps when available.
        try:
            with open(tmp_path, "rb") as audio_file:
                transcript = client.audio.transcriptions.create(
                    model=settings.GROQ_WHISPER_MODEL,
                    file=audio_file,
                    response_format="verbose_json",
                )
        finally:
            os.unlink(tmp_path)

        full_text = transcript.text or ""
        # Groq Whisper returns segments as dicts; OpenAI returns objects.
        # We normalise to dicts for consistent access.
        raw_segments = getattr(transcript, "segments", None) or []
        segments = []
        for s in raw_segments:
            if isinstance(s, dict):
                segments.append(s)
            else:
                # Pydantic/object-style — convert to dict
                segments.append({
                    "text": getattr(s, "text", ""),
                    "start": getattr(s, "start", 0.0),
                    "end": getattr(s, "end", 0.0),
                })
        duration = getattr(transcript, "duration", None) or 0.0

        # ── Step 3: Merge segments into larger, coherent chunks ──────────────
        # Each Whisper segment is ~5-15s. Merging several into ~300-token chunks
        # gives the LLM better context (no mid-sentence cuts) and reduces the
        # number of embeddings by ~10x.
        max_tokens = getattr(settings, "AUDIO_CHUNK_MAX_TOKENS", 300)
        chunks: list[Dict] = []
        buf_text = []
        buf_tokens = 0
        seg_start: float = 0.0
        seg_end: float = 0.0

        for i, seg in enumerate(segments):
            seg_text = (seg.get("text") or "").strip()
            if not seg_text:
                continue
            seg_tok = len(seg_text) // 4  # rough 1 token ≈ 4 chars

            if i == 0:
                seg_start = seg.get("start", 0.0)

            if buf_tokens + seg_tok > max_tokens and buf_text:
                # Flush current buffer as one chunk
                chunks.append({
                    "content": " ".join(buf_text),
                    "chunk_type": "transcript_segment",
                    "timestamp_start": seg_start,
                    "timestamp_end": seg_end,
                })
                buf_text = []
                buf_tokens = 0
                seg_start = seg.get("start", 0.0)

            buf_text.append(seg_text)
            buf_tokens += seg_tok
            seg_end = seg.get("end", 0.0)

        # Flush remaining buffer
        if buf_text:
            chunks.append({
                "content": " ".join(buf_text),
                "chunk_type": "transcript_segment",
                "timestamp_start": seg_start,
                "timestamp_end": seg_end,
            })

        # If Whisper returned no segments (rare edge case), fall back to full text
        if not chunks and full_text:
            chunks = [{"content": full_text, "chunk_type": "transcript_segment",
                       "timestamp_start": 0.0, "timestamp_end": float(duration)}]

        # ── Step 4: Compress full transcript for metadata search ─────────────
        # Stored in family_data so keyword/metadata search always works without
        # needing a vector lookup. zlib compression gives ~70% size reduction.
        compressed_transcript = ""
        if full_text:
            compressed_bytes = zlib.compress(full_text.encode("utf-8"), level=6)
            compressed_transcript = base64.b64encode(compressed_bytes).decode("ascii")

        # ── Step 5: Extract top-10 keywords (simple TF-IDF approximation) ───
        keywords = _extract_keywords(full_text, top_n=10)

        # ── Step 6: Optionally delete raw audio from MinIO ───────────────────
        # Enabled when KEEP_AUDIO_RAW=False (default).
        # The storage key used to upload the audio is in file_record.storage_key
        # but we don't have it here at parse time — the deletion is handled in
        # tasks.py after _save_chunks() succeeds by checking settings.KEEP_AUDIO_RAW.
        # We signal this via family_data["audio_raw_pending_deletion"] = True.

        family_data: Dict = {
            "duration_seconds": duration,
            "word_count": len(full_text.split()),
            "segment_count": len(segments),
            "merged_chunk_count": len(chunks),
            "language": transcript.language or "en",
            "keywords": keywords,
            # Compressed transcript — decompress with:
            #   import zlib, base64
            #   zlib.decompress(base64.b64decode(family_data["transcript_compressed"]))
            "transcript_compressed": compressed_transcript,
            "audio_raw_pending_deletion": not getattr(settings, "KEEP_AUDIO_RAW", False),
        }

        return {
            "title": filename,
            "language": transcript.language or "en",
            "family_data": family_data,
            "chunks": chunks,
            "entities": _extract_entities(full_text),
        }

    except Exception as e:
        return _error_result(filename, str(e))


# ─── CAD ─────────────────────────────────────────────────────────────────────

def _parse_cad(content: bytes, filename: str, mime_type: str) -> Dict:
    """
    MVP: Extract metadata from DXF/DWG title blocks via text extraction.
    Full geometry parsing deferred to post-MVP.
    """
    try:
        ext = filename.lower().rsplit(".", 1)[-1]
        text_content = ""
        family_data = {"format": ext}

        if ext == "dxf":
            import ezdxf
            doc = ezdxf.read(io.StringIO(content.decode("utf-8", errors="replace")))
            msp = doc.modelspace()
            texts = [e.dxf.text for e in msp if e.dxftype() in ("TEXT", "MTEXT") and hasattr(e.dxf, "text")]
            text_content = "\n".join(texts)
            layers = [layer.dxf.name for layer in doc.layers]
            family_data["layers"] = layers

        # For DWG and other formats: extract any embedded text
        if not text_content:
            # Naive text extraction from binary
            text_content = re.sub(rb"[^\x20-\x7E\n]", b" ", content).decode("ascii", errors="ignore")
            text_content = " ".join(text_content.split())[:5000]

        # Extract drawing number, project info from title block patterns
        drawing_num = re.search(r"DWG\.?\s*NO\.?\s*[:\s]*([A-Z0-9\-]+)", text_content, re.I)
        if drawing_num:
            family_data["drawing_number"] = drawing_num.group(1)

        chunks = []
        if text_content:
            chunks.append({"content": text_content[:4000], "chunk_type": "drawing_text"})

        return {
            "title": family_data.get("drawing_number", filename),
            "language": "en",
            "family_data": family_data,
            "chunks": chunks,
            "entities": _extract_entities(text_content),
        }
    except Exception as e:
        return _error_result(filename, str(e))


# ─── Operational / Exports ────────────────────────────────────────────────────

def _parse_operational(content: bytes, filename: str, mime_type: str) -> Dict:
    try:
        ext = filename.lower().rsplit(".", 1)[-1]

        if ext == "json":
            import json
            data = json.loads(content)
            text = json.dumps(data, indent=2)[:10000]
            return {
                "title": filename,
                "language": "en",
                "family_data": {"format": "json", "top_keys": list(data.keys())[:20] if isinstance(data, dict) else []},
                "chunks": [{"content": text[i:i+1000], "chunk_type": "json_block"} for i in range(0, len(text), 1000)],
                "entities": _extract_entities(text),
            }

        elif ext == "xml":
            text = content.decode("utf-8", errors="replace")[:10000]
            clean = re.sub(r"<[^>]+>", " ", text)
            return {
                "title": filename,
                "language": "en",
                "family_data": {"format": "xml"},
                "chunks": [{"content": clean[i:i+1000], "chunk_type": "xml_block"} for i in range(0, min(len(clean), 8000), 1000)],
                "entities": _extract_entities(clean),
            }

        else:
            return _parse_plaintext(content, filename)

    except Exception as e:
        return _error_result(filename, str(e))


# ─── Unknown ─────────────────────────────────────────────────────────────────

def _parse_unknown(content: bytes, filename: str, mime_type: str) -> Dict:
    text = re.sub(rb"[^\x20-\x7E\n]", b" ", content).decode("ascii", errors="ignore")
    text = " ".join(text.split())[:5000]
    return {
        "title": filename,
        "language": "en",
        "family_data": {},
        "chunks": [{"content": text, "chunk_type": "raw_text"}] if text else [],
        "entities": _extract_entities(text),
    }


# ─── Keyword extraction ────────────────────────────────────────────────────────

# Common stop-words to filter out of keyword extraction
_STOP_WORDS = {
    "the", "a", "an", "and", "or", "but", "in", "on", "at", "to", "for",
    "of", "is", "are", "was", "were", "be", "been", "being", "have", "has",
    "had", "do", "does", "did", "will", "would", "could", "should", "may",
    "might", "that", "this", "these", "those", "it", "its", "with", "from",
    "by", "as", "up", "out", "not", "we", "i", "you", "he", "she", "they",
    "our", "your", "their", "my", "his", "her", "all", "so", "can", "if",
}


def _extract_keywords(text: str, top_n: int = 10) -> list:
    """
    Simple TF-based keyword extraction — no external dependencies.
    Filters stop-words and short tokens, returns the most frequent terms.
    Adequate for generating searchable keywords from transcripts.
    """
    from collections import Counter
    if not text:
        return []
    words = re.findall(r'\b[a-zA-Z]{4,}\b', text.lower())
    filtered = [w for w in words if w not in _STOP_WORDS]
    counter = Counter(filtered)
    return [word for word, _ in counter.most_common(top_n)]


# ─── Entity extraction ────────────────────────────────────────────────────────

EMAIL_PATTERN = re.compile(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,7}\b')
PHONE_PATTERN = re.compile(r'\b(?:\+?\d{1,3}[-.\s]?)?\(?\d{3}\)?[-.\s]?\d{3}[-.\s]?\d{4}\b')
PROPER_NOUN_PATTERN = re.compile(r'\b([A-Z][a-zA-Z]+(?:\s+[A-Z][a-zA-Z]+){1,3})\b')
MEASUREMENT_PATTERN = re.compile(r'\b(\d+(?:\.\d+)?\s*(?:psi|bar|deg|c|f|rpm|v|kw|hz|m|cm|mm|kg|l|gpm|cfm|psig|kpa))\b', re.IGNORECASE)

def _extract_entities(text: str, is_image: bool = False) -> List[Dict]:
    entities = []
    seen = set()

    def add_entity(etype: str, val: str, conf: float = 0.90):
        v = val.strip()
        if len(v) > 2 and v.lower() not in seen:
            seen.add(v.lower())
            entities.append({"type": etype, "value": v, "confidence": conf})

    # Equipment IDs
    for match in EQUIPMENT_PATTERN.finditer(text):
        add_entity("equipment", match.group(1), 0.95)

    # Standards
    for match in STANDARD_PATTERN.finditer(text):
        add_entity("standard", match.group(0), 0.90)

    # Measurements / Engineering parameters (especially common in P&IDs / handwriting)
    for match in MEASUREMENT_PATTERN.finditer(text):
        add_entity("measurement", match.group(1), 0.88)

    # Emails
    for match in EMAIL_PATTERN.finditer(text):
        add_entity("email", match.group(0), 0.95)

    # Phones
    for match in PHONE_PATTERN.finditer(text):
        add_entity("phone", match.group(0), 0.85)

    # Proper Nouns / Orgs / Names
    for match in PROPER_NOUN_PATTERN.finditer(text):
        val = match.group(1)
        if len(val.split()) <= 4 and not any(w.lower() in _STOP_WORDS for w in val.split()):
            add_entity("proper_noun", val, 0.80)

    # Keep high number of entities across both small and large documents for rich knowledge graph queries
    max_entities = 150
    return entities[:max_entities]


# ─── Utilities ────────────────────────────────────────────────────────────────

def _split_into_chunks(text: str, max_chars: int = 2000, overlap: int = 200) -> List[str]:
    """Split text into overlapping chunks."""
    if not text or len(text) <= max_chars:
        return [text] if text else []
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + max_chars, len(text))
        # Try to break at sentence boundary
        if end < len(text):
            last_period = text.rfind(".", start, end)
            if last_period > start + max_chars // 2:
                end = last_period + 1
        chunks.append(text[start:end].strip())
        if end >= len(text):
            break
        start = max(start + 1, end - overlap)
    return [c for c in chunks if c.strip()]


def _merge_small_chunks(chunks: List[Dict], min_chars: int = 100) -> List[Dict]:
    """Merge very small chunks into the previous one."""
    merged = []
    for chunk in chunks:
        if merged and len(chunk["content"]) < min_chars:
            merged[-1]["content"] += " " + chunk["content"]
        else:
            merged.append(chunk)
    return merged


def _group_into_sections(chunks: List[Dict], max_chars: int = 2200) -> List[Dict]:
    """
    Group fine-grained paragraph/page/row chunks into larger section-wise chunks.
    This reduces total chunk count by ~6x to ~10x, enabling rapid embedding across
    10,000+ page datasets while fitting perfectly into local embedding token windows.
    """
    if not chunks:
        return []
    sections = []
    curr_section = []
    curr_chars = 0
    first_meta = {}

    for c in chunks:
        content = c.get("content", "").strip()
        if not content:
            continue
        
        if curr_chars + len(content) > max_chars and curr_section:
            sections.append({
                "content": "\n\n".join(curr_section),
                **{k: v for k, v in first_meta.items() if k != "content"}
            })
            curr_section = [content]
            curr_chars = len(content)
            first_meta = {k: v for k, v in c.items() if k != "content"}
        else:
            if not curr_section:
                first_meta = {k: v for k, v in c.items() if k != "content"}
            curr_section.append(content)
            curr_chars += len(content) + 2

    if curr_section:
        sections.append({
            "content": "\n\n".join(curr_section),
            **{k: v for k, v in first_meta.items() if k != "content"}
        })

    return sections


def _guess_title(filename: str, text: str) -> str:
    """Try to extract title from first significant line of text."""
    lines = [l.strip() for l in text.split("\n") if l.strip() and len(l.strip()) > 5]
    if lines:
        first = lines[0][:120]
        if len(first) > 10:
            return first
    return filename


def _error_result(filename: str, error: str) -> Dict:
    return {
        "title": filename,
        "language": "en",
        "family_data": {"parse_error": error[:500]},
        "chunks": [],
        "entities": [],
    }
